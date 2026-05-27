import streamlit as st

st.set_page_config(
    page_title="DeepInsight | AI Analytics",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded",
)

import pandas as pd
import httpx
import plotly.graph_objects as go
import plotly.express as px
import plotly.io as pio
import json
import io
import time
from datetime import datetime

# ── Configuration ─────────────────────────────────────────────

API_BASE_URL = st.sidebar.text_input("API Base URL", value="http://localhost:8000")
AUTH_TOKEN = st.sidebar.text_input("Auth Token (Optional)", value="", type="password")

# ─────────────────────────────────────────────────────────────
# PPTX Builder — dark enterprise theme (no API call needed)
# ─────────────────────────────────────────────────────────────

def build_section_pptx(
    title: str,
    subtitle: str,
    metrics: dict,
    findings: str,
    recommendations: str,
    risks: str,
    conclusion: str,
    dataset_name: str = "Dataset",
    chart_fig=None,
) -> bytes:
    """Build a professional 6-slide dark PPTX from analysis results."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BG   = RGBColor(0x0A, 0x0D, 0x1A)
    CARD = RGBColor(0x0F, 0x17, 0x2A)
    IND  = RGBColor(0x63, 0x66, 0xF1)
    CYN  = RGBColor(0x06, 0xB6, 0xD4)
    PUR  = RGBColor(0x8B, 0x5C, 0xF6)
    WHT  = RGBColor(0xF8, 0xFA, 0xFC)
    MUT  = RGBColor(0x94, 0xA3, 0xB8)
    RED  = RGBColor(0xEF, 0x44, 0x44)
    GRN  = RGBColor(0x10, 0xB9, 0x81)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(slide, c=BG):
        f = slide.background.fill; f.solid(); f.fore_color.rgb = c

    def rect(slide, l, t, w, h, c):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

    def txt(slide, text, l, t, w, h, sz=14, bold=False, c=WHT, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = c; r.font.name = "Calibri"

    def line(slide, top, c=IND):
        s = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.33), Pt(2))
        s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

    # ── Slide 1: Cover ────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    rect(s, Inches(0), Inches(0), Inches(0.08), Inches(7.5), IND)
    txt(s, "DeepInsights AI", Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), 12, True, IND)
    txt(s, title, Inches(0.5), Inches(2.2), Inches(10), Inches(1.2), 36, True, WHT)
    txt(s, subtitle, Inches(0.5), Inches(3.5), Inches(10), Inches(0.8), 15, False, MUT)
    txt(s, f"Dataset: {dataset_name}", Inches(0.5), Inches(4.4), Inches(8), Inches(0.4), 12, False, MUT)
    txt(s, f"Generated: {datetime.utcnow().strftime('%B %d, %Y  %H:%M UTC')}", Inches(0.5), Inches(6.9), Inches(8), Inches(0.4), 11, False, MUT)
    rect(s, Inches(10), Inches(2), Inches(3), Inches(3.5), CARD)
    txt(s, "AI\nPowered\nAnalytics", Inches(10.1), Inches(2.8), Inches(2.8), Inches(2.5), 22, True, PUR, PP_ALIGN.CENTER)

    # ── Slide 2: Key Metrics ──────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    rect(s, Inches(0), Inches(0), Inches(0.08), Inches(7.5), CYN)
    txt(s, "Key Performance Metrics", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6), 26, True, WHT)
    line(s, Inches(1.0), CYN)
    cols = 4; metric_items = list(metrics.items())[:8]
    for i, (label, val) in enumerate(metric_items):
        col = i % cols; row = i // cols
        lft = Inches(0.4 + col * 3.2); top = Inches(1.3 + row * 2.0)
        rect(s, lft, top, Inches(2.9), Inches(1.6), CARD)
        txt(s, str(val), lft + Inches(0.15), top + Inches(0.15), Inches(2.6), Inches(0.8), 24, True, IND, PP_ALIGN.CENTER)
        txt(s, label, lft + Inches(0.15), top + Inches(1.0), Inches(2.6), Inches(0.4), 11, False, MUT, PP_ALIGN.CENTER)

    # ── Slide 3: Chart (if available) ─────────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    rect(s, Inches(0), Inches(0), Inches(0.08), Inches(7.5), PUR)
    txt(s, "Visual Analysis", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6), 26, True, WHT)
    line(s, Inches(1.0), PUR)
    if chart_fig is not None:
        try:
            img_bytes = pio.to_image(chart_fig, format="png", width=1100, height=500, scale=1.5)
            s.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5))
        except Exception:
            txt(s, "Chart rendering not available. View charts in the interactive dashboard.", Inches(0.5), Inches(3.5), Inches(12), Inches(0.8), 14, False, MUT, PP_ALIGN.CENTER)
    else:
        txt(s, "See the interactive dashboard for visualizations.", Inches(0.5), Inches(3.5), Inches(12), Inches(0.8), 14, False, MUT, PP_ALIGN.CENTER)

    # ── Slide 4: Key Findings ─────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    rect(s, Inches(0), Inches(0), Inches(0.08), Inches(7.5), GRN)
    txt(s, "Key Findings & Insights", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6), 26, True, WHT)
    line(s, Inches(1.0), GRN)
    txt(s, findings[:900] if findings else "No findings available.", Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), 13, False, MUT)

    # ── Slide 5: Risks & Recommendations ─────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    rect(s, Inches(0), Inches(0), Inches(0.08), Inches(7.5), RED)
    txt(s, "Risks & Strategic Recommendations", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6), 26, True, WHT)
    line(s, Inches(1.0), RED)
    half = Inches(6.0)
    txt(s, "⚠  Risk Alerts", Inches(0.5), Inches(1.1), half, Inches(0.4), 14, True, RED)
    txt(s, risks[:400] if risks else "No risks identified.", Inches(0.5), Inches(1.6), half, Inches(2.5), 12, False, MUT)
    txt(s, "✅  Recommendations", Inches(6.8), Inches(1.1), half, Inches(0.4), 14, True, GRN)
    txt(s, recommendations[:400] if recommendations else "Review findings above.", Inches(6.8), Inches(1.6), half, Inches(2.5), 12, False, MUT)
    line(s, Inches(4.4))
    txt(s, conclusion[:250] if conclusion else "", Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.5), 13, False, RGBColor(0xA5, 0xB4, 0xFC))

    # ── Slide 6: Closing ──────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    rect(s, Inches(0), Inches(0), Inches(0.08), Inches(7.5), IND)
    txt(s, "DeepInsights", Inches(2), Inches(2.5), Inches(9.33), Inches(1.2), 48, True, WHT, PP_ALIGN.CENTER)
    txt(s, "Enterprise AI Analytics Platform", Inches(2), Inches(3.7), Inches(9.33), Inches(0.6), 18, False, IND, PP_ALIGN.CENTER)
    line(s, Inches(4.5))
    txt(s, "Generated by DeepInsights AI  •  Confidential", Inches(2), Inches(4.8), Inches(9.33), Inches(0.5), 12, False, MUT, PP_ALIGN.CENTER)

    out = io.BytesIO(); prs.save(out); return out.getvalue()


# ─────────────────────────────────────────────────────────────
# Executive Summary Caller
# ─────────────────────────────────────────────────────────────

def call_executive_summary(context: str, section_type: str, dataset_name: str = "Dataset") -> dict:
    """Call backend /api/ai/summarize and return sections dict."""
    try:
        resp = httpx.post(
            f"{API_BASE_URL}/api/ai/summarize",
            json={"context": context, "section_type": section_type, "dataset_name": dataset_name},
            headers={"Authorization": f"Bearer {AUTH_TOKEN}"} if AUTH_TOKEN else {},
            timeout=45.0,
        )
        if resp.status_code == 200:
            return resp.json()
    except Exception:
        pass
    return {}


# ─────────────────────────────────────────────────────────────
# Export Panel — rendered at the bottom of each results section
# ─────────────────────────────────────────────────────────────

def render_export_panel(
    section_type: str,    # "ml" | "forecast" | "anomaly"
    section_label: str,   # display name
    metrics: dict,
    context_str: str,
    chart_fig=None,
    state_key: str = "",
):
    """Renders the Export PPTX + Executive Summary panel."""
    dataset_name = (st.session_state.dataset_meta or {}).get("file_name", "Dataset")
    st.markdown("---")
    st.markdown(
        f'<div style="background:linear-gradient(135deg,rgba(99,102,241,0.08),rgba(139,92,246,0.08));'
        f'border:1px solid rgba(99,102,241,0.2);border-radius:14px;padding:1.2rem 1.5rem;margin-top:0.5rem">'
        f'<span style="font-size:1rem;font-weight:700;color:#a5b4fc">📤 Export & Insights</span>'
        f'<span style="font-size:0.8rem;color:#64748b;margin-left:8px">— {section_label}</span></div>',
        unsafe_allow_html=True,
    )
    st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)

    ec1, ec2 = st.columns(2)

    # ── PPTX Export ───────────────────────────────────────────
    with ec1:
        summ_key = f"exec_summ_{state_key}"
        summ = st.session_state.get(summ_key, {})
        if st.button(f"📊 Export PPTX Report", key=f"pptx_{state_key}", use_container_width=True):
            with st.spinner("Building enterprise PPTX…"):
                pptx_bytes = build_section_pptx(
                    title=f"{section_label} Report",
                    subtitle=f"AI-Generated Analysis Report — {dataset_name}",
                    metrics=metrics,
                    findings=summ.get("key_findings", context_str[:600]),
                    recommendations=summ.get("recommendations", "See full analysis for recommendations."),
                    risks=summ.get("risks_and_alerts", "Review detailed metrics for risk signals."),
                    conclusion=summ.get("conclusion", "Further analysis recommended."),
                    dataset_name=dataset_name,
                    chart_fig=chart_fig,
                )
                st.download_button(
                    label="⬇ Download PPTX",
                    data=pptx_bytes,
                    file_name=f"DeepInsights_{section_type}_{datetime.utcnow().strftime('%Y%m%d_%H%M')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"dl_{state_key}",
                    use_container_width=True,
                )

    # ── Executive Summary ─────────────────────────────────────
    with ec2:
        summ_key = f"exec_summ_{state_key}"
        if st.button(f"✨ Generate Executive Summary", key=f"summ_{state_key}", use_container_width=True):
            with st.spinner("AI is generating your executive summary…"):
                result = call_executive_summary(context_str, section_type, dataset_name)
                if result:
                    st.session_state[summ_key] = result
                else:
                    st.error("Summary generation failed. Check API connection.")

    # ── Show summary if already generated ─────────────────────
    summ = st.session_state.get(f"exec_summ_{state_key}", {})
    if summ:
        SECTION_ICONS = {
            "overview": ("🎯", "#6366f1"),
            "key_findings": ("🔍", "#06b6d4"),
            "performance_summary": ("📊", "#10b981"),
            "risks_and_alerts": ("⚠️", "#ef4444"),
            "recommendations": ("💡", "#f59e0b"),
            "conclusion": ("✅", "#8b5cf6"),
        }
        st.markdown(
            '<div style="background:rgba(10,13,26,0.9);border:1px solid rgba(99,102,241,0.25);'
            'border-radius:14px;padding:1.2rem 1.5rem;margin-top:0.8rem">',
            unsafe_allow_html=True,
        )
        st.markdown(f"#### 🤖 AI Executive Summary — {section_label}")
        for key, (icon, color) in SECTION_ICONS.items():
            val = summ.get(key, "")
            if val:
                label = key.replace("_", " ").title()
                with st.expander(f"{icon}  {label}", expanded=(key in ("overview", "key_findings"))):
                    st.markdown(
                        f'<div style="color:#cbd5e1;line-height:1.7;font-size:0.88rem">{val}</div>',
                        unsafe_allow_html=True,
                    )
        st.markdown("</div>", unsafe_allow_html=True)

        # Download summary as text
        summary_text = "\n\n".join(
            [f"## {k.replace('_',' ').title()}\n{v}" for k, v in summ.items() if v]
        )
        st.download_button(
            "📄 Download Summary (TXT)",
            data=summary_text.encode(),
            file_name=f"Executive_Summary_{section_type}_{datetime.utcnow().strftime('%Y%m%d')}.txt",
            mime="text/plain",
            key=f"dl_summ_{state_key}",
        )



# ── Custom CSS ────────────────────────────────────────────────

st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');
    
    .main { background: linear-gradient(160deg, #0a0f1a 0%, #0d1b2a 40%, #1b2838 100%); color: #e0e7ef; }
    
    .stButton>button {
        background: linear-gradient(135deg, #06b6d4, #8b5cf6) !important;
        color: white !important; border-radius: 10px; border: none;
        padding: 0.6rem 1.4rem; font-weight: 600;
        transition: all 0.3s ease; letter-spacing: 0.02em;
    }
    .stButton>button p, div[data-testid="stButton"] > button p {
        color: inherit !important;
    }
    .stButton>button:hover {
        background: linear-gradient(135deg, #22d3ee, #a78bfa) !important;
        box-shadow: 0 6px 20px rgba(6, 182, 212, 0.35); transform: translateY(-2px);
    }
    .stMetric {
        background: rgba(13, 27, 42, 0.8); padding: 1.5rem;
        border-radius: 14px; border: 1px solid rgba(6, 182, 212, 0.2);
        backdrop-filter: blur(12px);
    }
    div[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #0d1b2a 0%, #1b2838 100%);
        border-right: 1px solid rgba(6, 182, 212, 0.15);
    }
    h1, h2, h3 { color: #f0f4f8; font-family: 'Inter', sans-serif; }
    
    .status-card {
        background: linear-gradient(135deg, rgba(6,182,212,0.08) 0%, rgba(139,92,246,0.08) 100%);
        padding: 2rem; border-radius: 16px;
        border-left: 4px solid #06b6d4; margin-bottom: 1.5rem;
        backdrop-filter: blur(10px); border: 1px solid rgba(6,182,212,0.15);
    }
    .metric-card {
        background: rgba(13,27,42,0.85); border-radius: 14px;
        padding: 1.4rem; text-align: center; margin-bottom: 1rem;
        border: 1px solid rgba(6,182,212,0.15); backdrop-filter: blur(10px);
        transition: transform 0.2s, box-shadow 0.2s;
    }
    .metric-card:hover { transform: translateY(-3px); box-shadow: 0 8px 25px rgba(6,182,212,0.15); }
    .metric-value { font-size: 1.8rem; font-weight: 800; background: linear-gradient(135deg, #06b6d4, #a78bfa); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }
    .metric-label { font-size: 0.75rem; color: #7c8da0; text-transform: uppercase; letter-spacing: 0.08em; margin-top: 0.3rem; }
    
    .hero-banner {
        background: linear-gradient(135deg, rgba(6,182,212,0.12), rgba(139,92,246,0.12));
        border: 1px solid rgba(6,182,212,0.2); border-radius: 20px;
        padding: 2.5rem; text-align: center; margin-bottom: 2rem;
    }
    .hero-banner h1 { font-size: 2.2rem; font-weight: 800; margin-bottom: 0.5rem;
        background: linear-gradient(135deg, #22d3ee, #a78bfa, #f472b6); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }
    .hero-banner p { color: #94a3b8; font-size: 1rem; }
    
    div[data-testid="stTabs"] button { color: #94a3b8 !important; font-weight: 600 !important; }
    div[data-testid="stTabs"] button[aria-selected="true"] { color: #22d3ee !important; border-bottom-color: #06b6d4 !important; }
    
    .stExpander { border: 1px solid rgba(6,182,212,0.15) !important; border-radius: 12px !important; background: rgba(13,27,42,0.6) !important; }
    </style>
""", unsafe_allow_html=True)

# ── Helpers ───────────────────────────────────────────────────

def get_headers():
    headers = {}
    if AUTH_TOKEN:
        headers["Authorization"] = f"Bearer {AUTH_TOKEN}"
    return headers

def check_health():
    try:
        response = httpx.get(f"{API_BASE_URL}/health", timeout=5.0)
        return response.status_code == 200
    except:
        return False

def show_api_error(prefix: str, resp):
    try:
        data = resp.json()
        if isinstance(data, dict) and "detail" in data:
            st.error(f"❌ {prefix}: {data['detail']}")
        else:
            st.error(f"❌ {prefix}: {resp.text}")
    except Exception:
        st.error(f"❌ {prefix}: {resp.text}")

# ── State Management ──────────────────────────────────────────

if "dataset_id" not in st.session_state:
    st.session_state["dataset_id"] = None
if "dataset_meta" not in st.session_state:
    st.session_state["dataset_meta"] = None

# ── Header ────────────────────────────────────────────────────

st.markdown("""
<div class="hero-banner">
    <h1>⚡ DeepInsight AI Studio</h1>
    <p>Next-Generation Data Analytics &amp; Machine Learning Platform</p>
</div>
""", unsafe_allow_html=True)

# ── Sidebar Status ────────────────────────────────────────────

st.sidebar.markdown("### ⚡ System")
is_healthy = check_health()
if is_healthy:
    st.sidebar.success("🟢 API Online")
else:
    st.sidebar.error("🔴 API Offline")
st.sidebar.markdown("---")
st.sidebar.markdown("**v2.0** · AI Studio")

# ── Main Content ──────────────────────────────────────────────

tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "📤 Upload", "📊 Overview", "🤖 AI Chat", "🧪 Model Comparison Lab", 
    "📈 Time Series Forecast", "🚨 Anomaly Detection"
])

# ── Tab 1: Upload ─────────────────────────────────────────────

with tab1:
    st.markdown('<div class="status-card"><h2>Upload New Dataset</h2><p>Upload a CSV, XLSX, or JSON file to begin your AI-powered analysis.</p></div>', unsafe_allow_html=True)
    
    uploaded_file = st.file_uploader("Choose a file", type=["csv", "xlsx", "json"])
    
    if uploaded_file is not None:
        if st.button("Start Processing"):
            with st.spinner("Analyzing data and uploading to storage..."):
                files = {"file": (uploaded_file.name, uploaded_file.getvalue(), uploaded_file.type)}
                try:
                    response = httpx.post(
                        f"{API_BASE_URL}/api/datasets/upload",
                        headers=get_headers(),
                        files=files,
                        timeout=60.0
                    )
                    
                    if response.status_code == 200:
                        meta = response.json()
                        st.session_state["dataset_meta"] = meta
                        st.session_state["dataset_id"] = meta["id"]
                        st.success(f"Successfully processed {uploaded_file.name}!")
                        st.balloons()
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error(f"Upload failed: {response.text}")
                except Exception as e:
                    st.error(f"Error connecting to API: {str(e)}")

# ── Tab 2: Overview ───────────────────────────────────────────

with tab2:
    st.markdown("### 📊 Dataset Overview")
    
    # Dataset Selector
    try:
        resp_list = httpx.get(f"{API_BASE_URL}/api/datasets/", headers=get_headers())
        if resp_list.status_code == 200:
            datasets = resp_list.json()
            if datasets:
                ds_options = {f"{d['file_name']} ({d['id'][:8]})": d['id'] for d in datasets}
                # Default index logic: find current ID in list
                current_id = st.session_state.get("dataset_id")
                opt_list = ["-- Select Dataset --"] + list(ds_options.keys())
                
                selected_ds_label = st.selectbox(
                    "Load an existing dataset:",
                    options=opt_list,
                    index=0 if not current_id else (opt_list.index(next((k for k, v in ds_options.items() if v == current_id), opt_list[0])) if any(v == current_id for v in ds_options.values()) else 0)
                )
                
                if selected_ds_label != "-- Select Dataset --":
                    selected_id = ds_options[selected_ds_label]
                    if st.session_state.get("dataset_id") != selected_id:
                        st.session_state["dataset_id"] = selected_id
                        # Fetch details
                        resp_det = httpx.get(f"{API_BASE_URL}/api/datasets/{selected_id}", headers=get_headers())
                        if resp_det.status_code == 200:
                            st.session_state["dataset_meta"] = resp_det.json()
                            st.rerun()
    except Exception as e:
        st.warning(f"Could not fetch dataset list: {e}")

    dataset_id = st.session_state.get("dataset_id")
    if not dataset_id:
        st.info("No dataset uploaded yet. Please go to the Upload tab.")
    else:
        meta = st.session_state.dataset_meta
        
        st.markdown(f"### Dataset: `{meta['file_name']}`")
        
        # Metrics Row
        m1, m2, m3, m4 = st.columns(4)
        m1.metric("Rows", f"{meta['row_count']:,}")
        m2.metric("Columns", f"{meta['column_count']}")
        m3.metric("Quality Score", f"{meta['quality_score']}%")
        m4.metric("Null %", f"{meta['null_percentage']}%")
        
        # Column Details
        st.markdown("---")
        st.markdown("### Column Metadata")
        cols_df = pd.DataFrame(meta["columns"])
        
        # Convert list columns to string for better display and Arrow compatibility
        if "sample_values" in cols_df.columns:
            cols_df["sample_values"] = cols_df["sample_values"].apply(lambda x: ", ".join(map(str, x)) if isinstance(x, list) else str(x))
            
        st.dataframe(cols_df, use_container_width=True)
        
        # Action Buttons
        st.markdown("### Run Advanced Analysis")
        
        analysis_options = {
            "Descriptive Statistics": "descriptive_stats",
            "Correlation Matrix": "correlation",
            "Distribution Analysis": "distribution",
            "Bar Charts (Categorical)": "bar",
            "Pie Charts (Categorical)": "pie",
        }
        
        selected_types = st.multiselect(
            "Select analysis types to generate:",
            options=list(analysis_options.keys()),
            default=["Descriptive Statistics", "Distribution Analysis"]
        )
        
        if st.button("Generate Selected Analysis"):
            if not selected_types:
                st.warning("Please select at least one analysis type.")
            else:
                with st.spinner("Generating insights and charts..."):
                    types_to_run = [analysis_options[t] for t in selected_types]
                    try:
                        resp = httpx.post(
                            f"{API_BASE_URL}/api/analysis/run/{st.session_state.dataset_id}",
                            headers=get_headers(),
                            json={"analysis_types": types_to_run},
                            timeout=120.0
                        )
                        
                        if resp.status_code == 200:
                            st.success("Analysis complete!")
                            analysis_data = resp.json()
                            
                            # Render results for each analysis
                            for analysis in analysis_data.get("analyses", []):
                                with st.expander(f"Analysis: {analysis['analysis_type'].replace('_', ' ').title()}", expanded=True):
                                    # Show insights
                                    if analysis["results"].get("insights"):
                                        st.markdown("#### Key Insights")
                                        for insight in analysis["results"]["insights"]:
                                            st.info(insight)
                                    
                                    # Render charts
                                    for chart in analysis.get("charts", []):
                                        st.markdown(f"#### {chart['title']}")
                                        try:
                                            # Plotly figure data is in chart['data']
                                            fig_dict = chart["data"]
                                            st.plotly_chart(fig_dict, use_container_width=True)
                                        except Exception as e:
                                            st.error(f"Could not render chart: {str(e)}")
                                    
                                    # Show detailed data in JSON
                                    if st.checkbox("Show Raw Data", key=f"raw_{analysis['analysis_type']}"):
                                        st.json(analysis["results"])
                        else:
                            show_api_error("Analysis failed", resp)
                    except Exception as e:
                        st.error(f"Error connecting to API: {str(e)}")

# ── Tab 3: AI Chat ────────────────────────────────────────────

# Chip CSS
st.markdown("""
<style>
/* ── Follow-up question chips — dark theme ── */
.fuq-label{font-size:11px;color:#475569;margin-bottom:6px;display:flex;align-items:center;gap:4px;}
div[data-testid="stButton"] > button {
    font-family:'Inter',system-ui,sans-serif !important;
}
/* Override chip buttons to look dark */
div[data-testid="column"] div[data-testid="stButton"] > button {
    background: rgba(15,23,42,0.95) !important;
    border: 1px solid rgba(99,102,241,0.30) !important;
    color: #a5b4fc !important;
    border-radius: 24px !important;
    padding: 6px 16px !important;
    font-size: 12px !important;
    font-weight: 500 !important;
    transition: all 0.2s !important;
    box-shadow: 0 0 0 0 rgba(99,102,241,0) !important;
    white-space: normal !important;
    word-wrap: break-word !important;
    height: auto !important;
    line-height: 1.4 !important;
    text-align: left !important;
}
div[data-testid="column"] div[data-testid="stButton"] > button:hover {
    background: rgba(30,41,59,0.98) !important;
    border-color: rgba(99,102,241,0.65) !important;
    color: #e0e7ff !important;
    box-shadow: 0 0 16px rgba(99,102,241,0.22) !important;
    transform: translateY(-1px) !important;
}
</style>
""", unsafe_allow_html=True)


def _chip_class(q):
    l = q.lower()
    if any(k in l for k in ["forecast","predict","future"]): return "🔵"
    if any(k in l for k in ["anomal","outlier","risk","unusual"]): return "🔴"
    if any(k in l for k in ["correlat","relation","between"]): return "🟣"
    if any(k in l for k in ["trend","growth","over time","change"]): return "🟢"
    if any(k in l for k in ["model","ml","classif","accuracy"]): return "🟡"
    return "🔷"


def fetch_smart_followups(session_id):
    try:
        resp = httpx.post(
            f"{API_BASE_URL}/api/chat/{session_id}/followups",
            headers=get_headers(), timeout=20.0
        )
        if resp.status_code == 200:
            return resp.json().get("questions", [])
    except Exception:
        pass
    return []


# Speech-to-Text HTML
SPEECH_TO_TEXT_HTML = """
<div id="stt">
  <button id="mb" onclick="toggle()">
    <svg id="mi" viewBox="0 0 24 24"><path d="M12 14c1.66 0 3-1.34 3-3V5c0-1.66-1.34-3-3-3S9 3.34 9 5v6c0 1.66 1.34 3 3 3z"/><path d="M17 11c0 2.76-2.24 5-5 5s-5-2.24-5-5H5c0 3.53 2.61 6.43 6 6.92V21h2v-3.08c3.39-.49 6-3.39 6-6.92h-2z"/></svg>
  </button>
  <span id="st">🎤 Click to speak</span>
  <span id="dots"></span>
  <span id="pv"></span>
</div>
<style>
*{margin:0;padding:0;box-sizing:border-box}
html,body{background:transparent;overflow:hidden;height:100%}
#stt{
  display:flex;align-items:center;gap:10px;
  padding:6px 14px;height:42px;
  background:rgba(15,23,42,0.9);
  border-radius:24px;border:1px solid #334155;
  font-family:Inter,-apple-system,system-ui,sans-serif;
  width:fit-content;margin:0 auto;
}
#mb{
  width:32px;height:32px;min-width:32px;
  border-radius:50%;border:2px solid #6366f1;
  background:rgba(99,102,241,0.1);
  cursor:pointer;display:flex;align-items:center;justify-content:center;
  transition:all 0.3s;
}
#mb:hover{background:rgba(99,102,241,0.25);transform:scale(1.1)}
#mb.on{background:rgba(239,68,68,0.15);border-color:#ef4444;animation:glow 1.5s infinite}
#mb svg{width:16px;height:16px;fill:#a5b4fc;transition:fill 0.3s}
#mb.on svg{fill:#ef4444}
@keyframes glow{
  0%{box-shadow:0 0 0 0 rgba(239,68,68,0.4)}
  70%{box-shadow:0 0 0 10px rgba(239,68,68,0)}
  100%{box-shadow:0 0 0 0 rgba(239,68,68,0)}
}
#st{font-size:12px;color:#94a3b8;white-space:nowrap}
#st.ok{color:#34d399}
#st.err{color:#f87171}
#st.on{color:#e2e8f0}
#dots{display:none}
#dots.on{display:inline-flex;gap:3px;align-items:center}
#dots .d{width:4px;height:4px;background:#ef4444;border-radius:50%;animation:b 1.4s infinite both}
#dots .d:nth-child(2){animation-delay:.16s}
#dots .d:nth-child(3){animation-delay:.32s}
@keyframes b{0%,80%,100%{transform:scale(0)}40%{transform:scale(1)}}
#pv{font-size:11px;color:#cbd5e1;max-width:180px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;display:none}
#pv.on{display:inline}
</style>
<script>
const SR=window.SpeechRecognition||window.webkitSpeechRecognition;
const mb=document.getElementById('mb'),st=document.getElementById('st');
const dots=document.getElementById('dots'),pv=document.getElementById('pv');
let rec,on=false,tid;

if(!SR){st.textContent='⚠ Not supported in this browser';st.className='err';mb.style.opacity='.4';mb.style.cursor='not-allowed'}

function inject(t){
  try{
    const ta=window.parent.document.querySelector('textarea[data-testid="stChatInputTextArea"]');
    if(!ta){st.textContent='⚠ Chat input not found';st.className='err';return}
    const s=Object.getOwnPropertyDescriptor(window.HTMLTextAreaElement.prototype,'value').set;
    s.call(ta,t);ta.dispatchEvent(new Event('input',{bubbles:true}));ta.focus();
    st.textContent='✅ Press Enter to send';st.className='ok';
    clearTimeout(tid);tid=setTimeout(()=>{st.textContent='🎤 Click to speak';st.className='';},3500);
  }catch(e){st.textContent='⚠ Insert failed';st.className='err'}
}

function toggle(){
  if(!SR)return;
  if(on){rec.stop();return}
  rec=new SR();rec.continuous=false;rec.interimResults=true;rec.lang='en-US';
  rec.onstart=()=>{
    on=true;mb.className='on';st.textContent='Listening...';st.className='on';
    dots.className='on';dots.innerHTML='<span class="d"></span><span class="d"></span><span class="d"></span>';
    pv.className='on';pv.textContent='';
  };
  rec.onresult=(e)=>{
    let f='',im='';
    for(let i=0;i<e.results.length;i++){if(e.results[i].isFinal)f+=e.results[i][0].transcript;else im+=e.results[i][0].transcript;}
    pv.textContent=f+im||'...';
    if(f){inject(f.trim());rec.stop();}
  };
  rec.onerror=(e)=>{on=false;mb.className='';dots.className='';pv.className='';const m={'not-allowed':'🚫 Denied','no-speech':'🔇 No speech','audio-capture':'⚠ No mic','network':'⚠ Network'};st.textContent=m[e.error]||('⚠ '+e.error);st.className='err';clearTimeout(tid);tid=setTimeout(()=>{st.textContent='🎤 Click to speak';st.className='';},3500);};
  rec.onend=()=>{on=false;mb.className='';dots.className='';pv.className='';};
  try{rec.start();}catch(x){st.textContent='⚠ Mic failed';st.className='err';}
}
</script>
"""

with tab3:
    import streamlit.components.v1 as components

    if not st.session_state.dataset_id:
        st.info("Upload a dataset first to start chatting.")
    else:
        # ── Init session state ─────────────────────────────────
        if "messages" not in st.session_state:
            st.session_state.messages = []
        if "follow_up_questions" not in st.session_state:
            st.session_state.follow_up_questions = []
        if "asked_questions" not in st.session_state:
            st.session_state.asked_questions = set()

        # ── Header ─────────────────────────────────────────────
        hcol1, hcol2 = st.columns([6, 1])
        with hcol1:
            st.markdown("### 💬 Chat with your Data")
            st.caption("Smart follow-up suggestions appear after each AI response.")
        with hcol2:
            st.markdown("<div style='padding-top:8px'>", unsafe_allow_html=True)
            components.html(SPEECH_TO_TEXT_HTML, height=52)
            st.markdown("</div>", unsafe_allow_html=True)

        # ── TTS helper ─────────────────────────────────────────
        def tts_button(text: str, key: str):
            safe = text.replace("\\", "\\\\").replace("`", "\\`")
            html = f"""
<div style="display:inline-block;margin-top:6px">
  <button id="tts_{key}"
    onclick="(function(){{
      const btn=document.getElementById('tts_{key}');
      if(window.__tts_{key}_speaking){{window.speechSynthesis.cancel();window.__tts_{key}_speaking=false;btn.innerHTML='🔊 Read aloud';btn.style.color='#a5b4fc';return;}}
      const u=new SpeechSynthesisUtterance(`{safe}`);
      u.rate=1;u.pitch=1;u.lang='en-US';
      u.onend=()=>{{window.__tts_{key}_speaking=false;btn.innerHTML='🔊 Read aloud';btn.style.color='#a5b4fc';}};
      window.speechSynthesis.cancel();window.__tts_{key}_speaking=true;
      btn.innerHTML='⏹ Stop';btn.style.color='#f87171';
      window.speechSynthesis.speak(u);
    }})()\"
    style="background:transparent;border:1px solid #6366f1;color:#a5b4fc;border-radius:20px;padding:3px 12px;font-size:12px;cursor:pointer;white-space:nowrap">
    🔊 Read aloud
  </button>
</div>"""
            components.html(html, height=38)

        # ── Render history ─────────────────────────────────────
        for i, message in enumerate(st.session_state.messages):
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
                if message["role"] == "assistant":
                    tts_button(message["content"], f"hist_{i}")

        # ── Follow-up chips ────────────────────────────────────
        pending_prompt = None
        fuq = [q for q in st.session_state.follow_up_questions
               if q not in st.session_state.asked_questions]

        if fuq and st.session_state.messages:
            st.markdown("**✨ Suggested follow-ups:**")
            row1 = fuq[:3]
            row2 = fuq[3:6]
            cols1 = st.columns(len(row1))
            for ci, q in enumerate(row1):
                icon = _chip_class(q)
                with cols1[ci]:
                    if st.button(f"{icon} {q}", key=f"chip_{ci}_{hash(q)}", use_container_width=True):
                        pending_prompt = q
            if row2:
                cols2 = st.columns(len(row2))
                for ci, q in enumerate(row2):
                    icon = _chip_class(q)
                    with cols2[ci]:
                        if st.button(f"{icon} {q}", key=f"chip2_{ci}_{hash(q)}", use_container_width=True):
                            pending_prompt = q

        # ── Input ──────────────────────────────────────────────
        typed = st.chat_input("Ask anything about your data…")
        prompt = pending_prompt or typed

        if prompt:
            st.session_state.asked_questions.add(prompt)
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)

            with st.chat_message("assistant"):
                msg_ph = st.empty()
                msg_ph.markdown("⏳ Thinking…")
                try:
                    if "chat_session_id" not in st.session_state:
                        s_resp = httpx.post(
                            f"{API_BASE_URL}/api/chat/sessions",
                            headers=get_headers(),
                            json={"dataset_id": st.session_state.dataset_id},
                            timeout=30.0
                        )
                        if s_resp.status_code == 200:
                            st.session_state.chat_session_id = s_resp.json()["id"]

                    if "chat_session_id" in st.session_state:
                        chat_resp = httpx.post(
                            f"{API_BASE_URL}/api/chat/{st.session_state.chat_session_id}/message",
                            headers=get_headers(),
                            json={"message": prompt},
                            timeout=60.0
                        )
                        if chat_resp.status_code == 200:
                            data = chat_resp.json()
                            full_response = data["content"]
                            api_fuqs = data.get("follow_up_questions", [])
                            msg_ph.markdown(full_response)
                            st.session_state.messages.append(
                                {"role": "assistant", "content": full_response}
                            )
                            tts_button(full_response, f"new_{len(st.session_state.messages)}")
                            # Use API follow-ups or fetch from dedicated endpoint
                            if not api_fuqs:
                                api_fuqs = fetch_smart_followups(
                                    st.session_state.chat_session_id
                                )
                            st.session_state.follow_up_questions = api_fuqs
                            st.rerun()
                        else:
                            msg_ph.markdown("❌ Could not get response from AI.")
                    else:
                        msg_ph.markdown("❌ Could not create chat session.")
                except Exception as e:
                    msg_ph.markdown(f"❌ Error: {str(e)}")



# ── Tab 4: ML Studio ──────────────────────────────────────────

with tab4:
    if not st.session_state.dataset_id:
        st.info("Upload a dataset first to train models.")
    else:
        st.header("🔬 ML Studio — Modern AutoML Dashboard")
        
        ml_tab1, ml_tab2, ml_tab3 = st.tabs(["🚀 Training Studio", "📈 Deep Visualizations", "🏆 Model Leaderboard"])
        
        with ml_tab1:
            st.markdown("### Model Configuration")
            
            # 1. Task & Model Selection
            meta = st.session_state.dataset_meta
            col_names = [c["name"] for c in meta["columns"]]
            
            c1, c2, c3 = st.columns(3)
            with c1:
                task_type = st.selectbox("Task Type", ["classification", "regression"])
            with c2:
                # Fetch available models based on task
                available_models = []
                try:
                    m_resp = httpx.get(f"{API_BASE_URL}/api/ml/available-models?problem_type={task_type}", headers=get_headers())
                    if m_resp.status_code == 200:
                        available_models = m_resp.json()
                except: pass
                model_name = st.selectbox("Choose Algorithm", available_models if available_models else ["Random Forest", "Linear Regression"])
            with c3:
                target_col = st.selectbox("Target Column", options=col_names)
                
            if st.button("🏗️ Train Model", use_container_width=True):
                with st.spinner(f"Training {model_name}..."):
                    try:
                        resp = httpx.post(
                            f"{API_BASE_URL}/api/ml/run-task",
                            headers=get_headers(),
                            params={
                                "dataset_id": st.session_state.dataset_id,
                                "task_type": task_type,
                                "model_name": model_name,
                                "target_column": target_col
                            },
                            timeout=180.0
                        )
                        
                        if resp.status_code == 200:
                            st.session_state["last_trained_full"] = resp.json()
                            st.success(f"Model {model_name} trained successfully!")
                        else:
                            st.error(f"Training failed: {resp.text}")
                    except Exception as e:
                        st.error(f"Error: {e}")
            
            # Display Metrics
            if "last_trained_full" in st.session_state:
                res = st.session_state["last_trained_full"]
                st.markdown("---")
                st.markdown(f"#### 📊 Performance: {res['model_name']}")
                
                metrics = res["metrics"]
                duration = res.get("duration", 0)
                
                m1, m2, m3, m4 = st.columns(4)
                if task_type == "classification":
                    accuracy = metrics.get("accuracy") if metrics.get("accuracy") is not None else 0
                    precision = metrics.get("precision") if metrics.get("precision") is not None else 0
                    recall = metrics.get("recall") if metrics.get("recall") is not None else 0
                    f1_score = metrics.get("f1_score") if metrics.get("f1_score") is not None else 0
                    with m1:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">Accuracy</div><div class="metric-value">{accuracy*100:.1f}%</div></div>', unsafe_allow_html=True)
                    with m2:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">Precision</div><div class="metric-value">{precision*100:.1f}%</div></div>', unsafe_allow_html=True)
                    with m3:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">Recall</div><div class="metric-value">{recall*100:.1f}%</div></div>', unsafe_allow_html=True)
                    with m4:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">F1 Score</div><div class="metric-value">{f1_score*100:.1f}%</div></div>', unsafe_allow_html=True)
                else:
                    r2_score = metrics.get("r2_score") if metrics.get("r2_score") is not None else 0
                    rmse = metrics.get("rmse") if metrics.get("rmse") is not None else 0
                    mae = metrics.get("mae") if metrics.get("mae") is not None else 0
                    with m1:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">R² Score</div><div class="metric-value">{r2_score:.4f}</div></div>', unsafe_allow_html=True)
                    with m2:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">RMSE</div><div class="metric-value">{rmse:.2f}</div></div>', unsafe_allow_html=True)
                    with m3:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">MAE</div><div class="metric-value">{mae:.2f}</div></div>', unsafe_allow_html=True)
                    with m4:
                        st.markdown(f'<div class="metric-card"><div class="metric-label">Time (s)</div><div class="metric-value">{duration:.2f}s</div></div>', unsafe_allow_html=True)

        with ml_tab2:
            if "last_trained_full" in st.session_state:
                res = st.session_state["last_trained_full"]
                charts = res.get("charts", {})
                
                c1, c2 = st.columns(2)
                with c1:
                    if "confusion_matrix" in charts:
                        st.plotly_chart(go.Figure(charts["confusion_matrix"]), use_container_width=True)
                    elif "residual_plot" in charts:
                        st.plotly_chart(go.Figure(charts["residual_plot"]), use_container_width=True)
                
                with c2:
                    if "feature_importance" in charts:
                        st.plotly_chart(go.Figure(charts["feature_importance"]), use_container_width=True)
                
                st.markdown("#### 🔍 Prediction Summary")
                st.info("This model is now ready for live inference in the 'AI Chat' or via the Prediction API.")
            else:
                st.info("Train a model to see deep visualizations.")

        with ml_tab3:
            st.markdown("### 🏆 Comparison Leaderboard")
            if "last_trained_full" in st.session_state:
                comparison = st.session_state["last_trained_full"].get("comparison", [])
                if comparison:
                    l_data = []
                    for m in comparison:
                        row = {
                            "Model": m["model_name"],
                            "Time (s)": round(m.get("training_time", 0) or 0, 2)
                        }
                        # Extract metrics
                        met_list = m.get("model_metrics") or []
                        met = (met_list[0] if met_list and met_list[0] is not None else {})
                        
                        if task_type == "classification":
                            accuracy = met.get('accuracy') if met.get('accuracy') is not None else 0
                            f1_score = met.get('f1_score') if met.get('f1_score') is not None else 0
                            row["Accuracy"] = f"{accuracy*100:.1f}%"
                            row["F1 Score"] = f"{f1_score*100:.1f}%"
                            sort_key = accuracy
                        else:
                            r2_score = met.get('r2_score') if met.get('r2_score') is not None else 0
                            rmse = met.get('rmse') if met.get('rmse') is not None else 0
                            row["R² Score"] = f"{r2_score:.4f}"
                            row["RMSE"] = f"{rmse:.2f}"
                            sort_key = r2_score
                        
                        row["_sort"] = sort_key
                        l_data.append(row)
                    
                    df_leader = pd.DataFrame(l_data).sort_values(by="_sort", ascending=False).reset_index(drop=True)
                    df_leader.index = df_leader.index + 1
                    df_leader.index.name = "Rank"
                    
                    st.dataframe(df_leader.drop(columns=["_sort"]), use_container_width=True)
                    
                    # Visual Comparison
                    st.markdown("#### 📊 Comparative Performance")
                    df_viz = pd.DataFrame(l_data).sort_values(by="_sort", ascending=True)
                    metric_name = "Accuracy" if task_type == "classification" else "R² Score"
                    
                    # Clean up percentage string for plotting if classification
                    if task_type == "classification":
                        df_viz["PlotValue"] = df_viz["Accuracy"].str.rstrip('%').astype('float') / 100
                    else:
                        df_viz["PlotValue"] = df_viz["R² Score"].astype('float')
                        
                    fig_comp = px.bar(
                        df_viz, x="PlotValue", y="Model", 
                        orientation="h", title=f"Model Rank by {metric_name}",
                        template="plotly_dark", color="PlotValue", color_continuous_scale="Viridis"
                    )
                    st.plotly_chart(fig_comp, use_container_width=True)

                    # ── Export Panel ──────────────────────────────────
                    res_ml = st.session_state.get("last_trained_full", {})
                    ml_metrics_export = {}
                    for row_m in l_data:
                        for k, v in row_m.items():
                            if k not in ("_sort", "Model"):
                                ml_metrics_export[f"{row_m['Model']} {k}"] = v
                    render_export_panel(
                        section_type="ml",
                        section_label="ML Model Comparison Lab",
                        metrics=ml_metrics_export,
                        context_str=json.dumps({
                            "model": res_ml.get("model_name"),
                            "metrics": res_ml.get("metrics", {}),
                            "comparison": [{"model": m["model_name"], "time": m.get("training_time")} for m in (res_ml.get("comparison") or [])],
                        }, default=str),
                        chart_fig=fig_comp,
                        state_key="ml_leaderboard",
                    )
                else:
                    st.info("No comparison data available yet.")
            else:
                st.info("Train at least one model to enable comparison.")

# ── Tab 5: Time Series Forecast ─────────────────────────────────

with tab5:
    if not st.session_state.dataset_id:
        st.info("Upload a dataset first to run time series forecasting.")
    else:
        st.header("📈 Time Series Forecasting")
        st.markdown("Predict future values using advanced statistical models.")
        
        meta = st.session_state.dataset_meta
        col_names = [c["name"] for c in meta["columns"]]
        date_candidates = [c["name"] for c in meta["columns"] if "date" in str(c.get("dtype", "")).lower() or "object" in str(c.get("dtype", "")).lower() or "datetime" in str(c.get("dtype", "")).lower()]
        numeric_cols = [c["name"] for c in meta["columns"] if "int" in str(c.get("dtype", "")).lower() or "float" in str(c.get("dtype", "")).lower()]

        fc_c1, fc_c2, fc_c3, fc_c4 = st.columns(4)
        with fc_c1:
            date_col = st.selectbox("Date Column (Optional)", ["Auto-detect"] + col_names, index=0)
        with fc_c2:
            val_col = st.selectbox("Target Value Column", ["Auto-detect"] + numeric_cols, index=0)
        with fc_c3:
            model_opts = {
                "ARIMA / SARIMA": "sarima",
                "Exponential Smoothing": "exponential_smoothing",
                "Moving Average": "moving_average",
                "Prophet": "prophet"
            }
            model_disp = st.selectbox("Forecasting Algorithm", list(model_opts.keys()))
            model_type = model_opts[model_disp]
        with fc_c4:
            periods = st.number_input("Forecast Horizon (periods)", min_value=1, max_value=365, value=30)
            
        if st.button("🚀 Run Forecast", use_container_width=True):
            with st.spinner(f"Running {model_disp} forecast..."):
                payload = {
                    "model": model_type,
                    "target_column": None if val_col == "Auto-detect" else val_col,
                    "date_column": None if date_col == "Auto-detect" else date_col,
                    "forecast_horizon": periods
                }
                try:
                    resp = httpx.post(
                        f"{API_BASE_URL}/api/forecast/run/{st.session_state.dataset_id}",
                        headers=get_headers(),
                        json=payload,
                        timeout=180.0
                    )
                    if resp.status_code == 200:
                        st.success("Forecast generated successfully!")
                        st.session_state["last_forecast"] = resp.json()
                    else:
                        show_api_error("Forecast failed", resp)
                except Exception as e:
                    st.error(f"Error: {e}")
                    
        if "last_forecast" in st.session_state:
            res = st.session_state["last_forecast"]
            st.markdown("---")
            
            if res.get("charts"):
                st.plotly_chart(res["charts"][0]["data"], use_container_width=True)
            
            st.markdown("### 📊 Forecast Summary")
            r_data = res.get("results", {})
            model_name_mapped = {
                "sarima": "ARIMA / SARIMA (Seasonal Autoregressive Integrated Moving Average)",
                "exponential_smoothing": "Holt-Winters Exponential Smoothing",
                "moving_average": "Simple Moving Average (Rolling Mean)",
                "prophet": "Prophet (Additive Regression Trend + Seasonality)"
            }.get(r_data.get("model_type", "").lower(), r_data.get("model_type", "Unknown"))
            
            col_s1, col_s2, col_s3 = st.columns(3)
            col_s1.metric("Forecasting Algorithm", r_data.get("model_type", "Unknown").upper())
            col_s2.metric("Target Value Column", r_data.get("value_column", "Unknown"))
            col_s3.metric("Horizon (periods)", f"{r_data.get('periods', periods)}")
            
            with st.expander("⚙️ Model Parameters & Details", expanded=True):
                st.markdown(f"**Forecasting Method:** {model_name_mapped}")
                st.markdown(f"**Data Frequency:** `{r_data.get('frequency', 'Unknown')}`")
                st.markdown(f"**Detected Seasonal Period:** `{r_data.get('seasonal_period', 1)}`")
                if r_data.get("model_type") == "sarima":
                    st.markdown("**Order Parameters:** `(p, d, q) = (1, 1, 1)`")
                elif r_data.get("model_type") == "exponential_smoothing":
                    st.markdown("**Smoothing Parameters:** `Trend: Additive | Seasonality: Additive`")
                elif r_data.get("model_type") == "moving_average":
                    st.markdown("**Window Size:** `7 periods (default)`")
                elif r_data.get("model_type") == "prophet":
                    st.markdown("**Growth Mode:** `Linear` | **Seasonalities:** `Yearly, Weekly`")
                
                if st.checkbox("Show Raw Forecast JSON Data"):
                    st.json(r_data)

            # ── Forecast Export Panel ─────────────────────────────
            fc_metrics = {
                "Model": r_data.get("model_type", "N/A").upper(),
                "Target Column": r_data.get("value_column", "N/A"),
                "Periods Forecast": str(r_data.get("periods", periods)),
                "Frequency": r_data.get("frequency", "N/A"),
                "Seasonal Period": str(r_data.get("seasonal_period", "N/A")),
            }
            fc_fig = None
            try:
                if res.get("charts"):
                    fc_fig = go.Figure(res["charts"][0]["data"])
            except Exception:
                pass
            render_export_panel(
                section_type="forecast",
                section_label="Time Series Forecast",
                metrics=fc_metrics,
                context_str=json.dumps(r_data, default=str),
                chart_fig=fc_fig,
                state_key="forecast_results",
            )

# ── Tab 6: Anomaly Detection ──────────────────────────────────

with tab6:
    if not st.session_state.dataset_id:
        st.info("Upload a dataset first to run anomaly detection.")
    else:
        st.header("🚨 Anomaly Detection")
        st.markdown("Detect outliers, unusual patterns, and data anomalies.")
        
        an_c1, an_c2 = st.columns(2)
        with an_c1:
            an_method = st.selectbox("Detection Algorithm", [
                "isolation_forest", "lof", "dbscan", "iqr", "zscore"
            ], format_func=lambda x: x.replace("_", " ").title() if x != "lof" and x != "iqr" else x.upper())
        with an_c2:
            if an_method in ["isolation_forest", "lof"]:
                an_thresh = st.slider("Contamination (Sensitivity)", min_value=0.01, max_value=0.5, value=0.05, step=0.01)
            elif an_method == "dbscan":
                an_thresh = st.number_input("EPS Distance", value=0.5, step=0.1)
            elif an_method == "zscore":
                an_thresh = st.slider("Z-Score Cutoff", min_value=1.0, max_value=5.0, value=3.0, step=0.1)
            else:
                an_thresh = st.slider("IQR Multiplier", min_value=1.0, max_value=5.0, value=1.5, step=0.1)
                
        if st.button("🔍 Detect Anomalies", use_container_width=True):
            with st.spinner(f"Running {an_method} anomaly detection..."):
                payload = {
                    "method": an_method,
                    "threshold": an_thresh
                }
                try:
                    resp = httpx.post(
                        f"{API_BASE_URL}/api/anomaly/run/{st.session_state.dataset_id}",
                        headers=get_headers(),
                        json=payload,
                        timeout=180.0
                    )
                    if resp.status_code == 200:
                        st.success("Anomalies detected successfully!")
                        st.session_state["last_anomaly"] = resp.json()
                    else:
                        show_api_error("Anomaly detection failed", resp)
                except Exception as e:
                    st.error(f"Error: {e}")
                    
        if "last_anomaly" in st.session_state:
            res = st.session_state["last_anomaly"]
            st.markdown("---")
            total = res.get("results", {}).get("total_anomalies", 0)
            
            st.markdown(f'<div class="status-card" style="border-left-color: #ef4444;"><h3 style="margin:0;color:#ef4444">{total} Total Anomalies Detected</h3></div>', unsafe_allow_html=True)
            
            charts = res.get("charts", [])
            if charts:
                cols = st.columns(len(charts) if len(charts) <= 2 else 2)
                for idx, chart in enumerate(charts):
                    with cols[idx % 2]:
                        st.plotly_chart(chart["data"], use_container_width=True)
                        
            with st.expander("View Detailed Anomaly Stats"):
                st.json(res.get("results", {}).get("anomalies_by_column", {}))

            # ── Anomaly Export Panel ──────────────────────────────
            an_res = res.get("results", {})
            an_metrics = {
                "Total Anomalies": str(an_res.get("total_anomalies", 0)),
                "Method": an_res.get("method", an_method).upper(),
                "Threshold": str(an_thresh),
                "Columns Scanned": str(len(an_res.get("anomalies_by_column", {}))),
            }
            for col_name, col_info in list(an_res.get("anomalies_by_column", {}).items())[:4]:
                if isinstance(col_info, dict):
                    an_metrics[f"{col_name} anomalies"] = str(col_info.get("count", 0))
            an_fig = None
            try:
                if res.get("charts"):
                    an_fig = go.Figure(res["charts"][0]["data"])
            except Exception:
                pass
            render_export_panel(
                section_type="anomaly",
                section_label="Anomaly Detection",
                metrics=an_metrics,
                context_str=json.dumps(an_res, default=str),
                chart_fig=an_fig,
                state_key="anomaly_results",
            )

# ── Footer ────────────────────────────────────────────────────

st.markdown("---")
st.markdown(
    '<p style="text-align: center; color: #4a5568; font-size: 0.85rem;">⚡ DeepInsight AI Studio v2.0 · Built with FastAPI & Streamlit · Powered by Multi-Model AI</p>',
    unsafe_allow_html=True
)
