"""
DeepInsight Starter Suite — Document Intelligence Engine.

Extracts text, tables, and structures from unstructured documents
like PDF, DOCX, and PPTX.
"""

import logging
import io
from typing import Optional, Dict, Any, List
import json

import pandas as pd
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)

def parse_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF document."""
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text_content = []
        for page in doc:
            text_content.append(page.get_text())
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to parse PDF: {e}")
        raise ValueError(f"Could not read PDF file: {str(e)}")

def parse_docx(file_bytes: bytes) -> str:
    """Extract text from a Word document."""
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        text_content = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to parse DOCX: {e}")
        raise ValueError(f"Could not read DOCX file: {str(e)}")

def parse_pptx(file_bytes: bytes) -> str:
    """Extract text from PowerPoint presentations."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"PPTX parsing failed: {e}")
        raise ValueError("Failed to parse PPTX file.")


def chunk_document(text: str, max_chunk_size: int = 4000) -> List[str]:
    """Split document into manageable chunks for LLM processing."""
    # simple word-based chunker
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0
    
    for word in words:
        # crude length approx
        current_length += len(word) + 1 
        if current_length > max_chunk_size:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_length = len(word) + 1
        else:
            current_chunk.append(word)
            
    if current_chunk:
        chunks.append(" ".join(current_chunk))
        
    return chunks


async def generate_ai_summary(chunks: List[str]) -> Dict[str, Any]:
    """Generate a comprehensive summary from document chunks using LLM."""
    from services import llm_client
    
    if not chunks:
        return {"summary": "Empty document", "entities": [], "risks": []}
        
    # Process just the first few chunks to save tokens in MVP
    sample_text = "\n...\n".join(chunks[:3])
    
    system_prompt = (
        "You are an expert business analyst. Read the following document excerpt and extract:\n"
        "1. A concise executive summary\n"
        "2. Key entities or topics mentioned\n"
        "3. Any potential risks or warnings highlighted\n"
        "4. Actionable recommendations\n"
        "Format as JSON with keys: 'summary', 'entities', 'risks', 'recommendations'."
    )
    
    try:
        resp = await llm_client.chat_completion(system_prompt, sample_text)
        
        # Try to parse the JSON response
        # Using a very crude approach for now, assuming the LLM respects JSON format
        clean_text = resp.answer.strip()
        if clean_text.startswith("```json"):
            clean_text = clean_text[7:-3]
        elif clean_text.startswith("```"):
            clean_text = clean_text[3:-3]
            
        return json.loads(clean_text)
    except Exception as e:
        logger.error(f"AI summarization failed: {e}")
        return {
            "summary": "AI summary could not be generated.",
            "entities": [],
            "risks": [],
            "recommendations": []
        }

def extract_document_text(file_bytes: bytes, filename: str) -> str:
    """Route to correct parser based on extension."""
    ext = filename.split(".")[-1].lower()
    if ext == "pdf":
        return parse_pdf(file_bytes)
    elif ext in ["doc", "docx"]:
        return parse_docx(file_bytes)
    elif ext in ["ppt", "pptx"]:
        return parse_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported document format: {ext}")
