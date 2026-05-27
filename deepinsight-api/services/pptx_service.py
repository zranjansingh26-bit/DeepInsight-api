"""
DeepInsight — Professional PPTX Report Builder.

Generates enterprise-grade dark-themed PowerPoint presentations
with AI narrative, dataset stats, and embedded Plotly charts.
"""

import io
import logging
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ── Color Palette ─────────────────────────────────────────────
BG_DARK       = RGBColor(0x0A, 0x0D, 0x1A)   # #0a0d1a — near black
BG_CARD       = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a — slate-900
ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xF1)   # #6366f1 — indigo-500
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # #8b5cf6 — violet-500
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)   # #06b6d4 — cyan-500
TEXT_WHITE    = RGBColor(0xF8, 0xFA, 0xFC)   # #f8fafc
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # #94a3b8 — slate-400
TEXT_ACCENT   = RGBColor(0xA5, 0xB4, 0xFC)   # #a5b4fc — indigo-300
BORDER        = RGBColor(0x33, 0x41, 0x55)   # #334155 — slate-700

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color: RGBColor):
    """Fill slide background with solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def _add_text(slide, text: str, left, top, width, height,
              font_size=18, bold=False, color=TEXT_WHITE,
              align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def _add_divider(slide, top, color=ACCENT_INDIGO):
    """Add a horizontal accent line."""
    line = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.33), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


class PptxBuilder:
    """Builds a professional dark-themed 9-slide PPTX presentation."""

    def __init__(self, report_data: dict):
        self.data = report_data
        self.dataset = report_data.get("dataset", {})
        self.analyses = report_data.get("analyses", [])
        self.narrative = report_data.get("executive_narrative", "")
        self.exec_summary = report_data.get("executive_summary", {})

        self.prs = Presentation()
        self.prs.slide_width  = SLIDE_W
        self.prs.slide_height = SLIDE_H

        # Use blank layout
        self.blank = self.prs.slide_layouts[6]

    # ── Slide 1: Cover ────────────────────────────────────────
    def add_cover(self):
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)

        # Left accent bar
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_INDIGO)

        # Large gradient-like box
        _add_rect(slide, Inches(0.08), Inches(2.2), Inches(8), Inches(3.5), BG_CARD)

        # Brand badge
        _add_text(slide, "DeepInsights AI", Inches(0.5), Inches(0.35),
                  Inches(6), Inches(0.6), font_size=13, bold=True, color=ACCENT_INDIGO)

        # Main title
        dataset_name = self.dataset.get("file_name", "Dataset Report")
        _add_text(slide, dataset_name, Inches(0.5), Inches(2.4),
                  Inches(9), Inches(1.0), font_size=36, bold=True, color=TEXT_WHITE)

        # Subtitle from narrative
        subtitle_lines = self.exec_summary.get("overview", "AI-Generated Executive Report")
        subtitle = subtitle_lines[:160] if len(subtitle_lines) > 160 else subtitle_lines
        _add_text(slide, subtitle, Inches(0.5), Inches(3.55),
                  Inches(9), Inches(1.2), font_size=14, color=TEXT_MUTED)

        # Date
        date_str = datetime.utcnow().strftime("%B %d, %Y")
        _add_text(slide, f"Generated: {date_str}", Inches(0.5), Inches(6.8),
                  Inches(5), Inches(0.4), font_size=11, color=TEXT_MUTED)

        # Right decorative accent
        _add_rect(slide, Inches(10), Inches(1.5), Inches(3), Inches(4.5), BG_CARD)
        _add_text(slide, "AI\nPowered\nAnalytics", Inches(10.3), Inches(2.8),
                  Inches(2.5), Inches(2.5), font_size=22, bold=True,
                  color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)

        logger.info("PPTX: Cover slide added")

    # ── Slide 2: Executive Summary ────────────────────────────
    def add_executive_summary(self):
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_INDIGO)

        _add_text(slide, "Executive Summary", Inches(0.5), Inches(0.3),
                  Inches(10), Inches(0.6), font_size=28, bold=True, color=TEXT_WHITE)
        _add_divider(slide, Inches(1.0))

        content = self.exec_summary.get("overview", self.narrative)
        if len(content) > 900:
            content = content[:900] + "..."

        _add_text(slide, content, Inches(0.5), Inches(1.2),
                  Inches(12.3), Inches(5.5), font_size=13, color=TEXT_MUTED)

    # ── Slide 3: Dataset Overview ─────────────────────────────
    def add_dataset_overview(self):
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_CYAN)

        _add_text(slide, "Dataset Overview", Inches(0.5), Inches(0.3),
                  Inches(10), Inches(0.6), font_size=28, bold=True, color=TEXT_WHITE)
        _add_divider(slide, Inches(1.0), color=ACCENT_CYAN)

        ds = self.dataset
        stats = [
            ("File", str(ds.get("file_name", "N/A"))),
            ("Type", str(ds.get("file_type", "N/A"))),
            ("Rows", f"{ds.get('row_count', 0):,}"),
            ("Columns", str(ds.get("column_count", 0))),
            ("Quality Score", f"{ds.get('quality_score', 0)}%"),
            ("Null %", f"{ds.get('null_percentage', 0):.1f}%"),
            ("Analyses Run", str(len(self.analyses))),
        ]

        cols = 4
        card_w = Inches(2.9)
        card_h = Inches(1.5)
        for i, (label, val) in enumerate(stats):
            col = i % cols
            row = i // cols
            left = Inches(0.4 + col * 3.15)
            top  = Inches(1.4 + row * 1.9)
            _add_rect(slide, left, top, card_w, card_h, BG_CARD)
            _add_text(slide, val, left + Inches(0.15), top + Inches(0.15),
                      card_w - Inches(0.3), Inches(0.75),
                      font_size=22, bold=True, color=ACCENT_INDIGO, align=PP_ALIGN.CENTER)
            _add_text(slide, label, left + Inches(0.15), top + Inches(0.9),
                      card_w - Inches(0.3), Inches(0.45),
                      font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ── Slides 4+: Analysis Result Slides ─────────────────────
    def add_analysis_slides(self):
        for a in self.analyses[:6]:  # cap at 6 analysis slides
            slide = self.prs.slides.add_slide(self.blank)
            _set_slide_bg(slide, BG_DARK)
            _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_PURPLE)

            title = a.get("type", "Analysis").replace("_", " ").title()
            _add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(10), Inches(0.6), font_size=26, bold=True, color=TEXT_WHITE)
            _add_divider(slide, Inches(1.0), color=ACCENT_PURPLE)

            # Try embed chart image first
            chart_added = self._try_embed_chart(slide, a)

            if not chart_added:
                # Fallback: show key metrics as text
                results = a.get("results", {})
                text_lines = []
                for k, v in list(results.items())[:12]:
                    if isinstance(v, (str, int, float)):
                        text_lines.append(f"• {k.replace('_',' ').title()}: {v}")
                content = "\n".join(text_lines) if text_lines else "Analysis results stored in system."
                _add_text(slide, content, Inches(0.5), Inches(1.3),
                          Inches(12.3), Inches(5.5), font_size=13, color=TEXT_MUTED)

    def _try_embed_chart(self, slide, analysis: dict) -> bool:
        """Attempt to render and embed a Plotly chart PNG."""
        try:
            import plotly.graph_objects as go
            import plotly.io as pio

            results = analysis.get("results", {})
            a_type = analysis.get("type", "")

            fig = None

            if a_type == "trend" and results.get("dates") and results.get("values"):
                fig = go.Figure()
                fig.add_trace(go.Scatter(
                    x=results["dates"], y=results["values"],
                    mode="lines+markers", name="Value",
                    line=dict(color="#6366f1", width=2),
                    marker=dict(color="#8b5cf6", size=4),
                ))
                if results.get("rolling_mean"):
                    fig.add_trace(go.Scatter(
                        x=results["dates"], y=results["rolling_mean"],
                        mode="lines", name="Rolling Mean",
                        line=dict(color="#06b6d4", width=1.5, dash="dot"),
                    ))

            elif a_type == "forecast" and results.get("forecast_dates"):
                fig = go.Figure()
                fig.add_trace(go.Scatter(
                    x=results.get("forecast_dates", []),
                    y=results.get("forecast_values", []),
                    mode="lines", name="Forecast",
                    line=dict(color="#6366f1", width=2),
                ))

            elif a_type in ("correlation", "descriptive_stats"):
                matrix = results.get("matrix", {})
                if matrix:
                    cols = list(matrix.keys())
                    z = [[matrix[c].get(r, 0) for r in cols] for c in cols]
                    fig = go.Figure(go.Heatmap(
                        z=z, x=cols, y=cols,
                        colorscale="Purples", showscale=False,
                    ))

            if fig is None:
                return False

            # Apply dark template
            fig.update_layout(
                paper_bgcolor="#0f172a",
                plot_bgcolor="#0f172a",
                font=dict(color="#94a3b8", size=10),
                margin=dict(l=20, r=20, t=30, b=20),
                height=400, width=900,
            )
            fig.update_xaxes(gridcolor="#1e293b")
            fig.update_yaxes(gridcolor="#1e293b")

            img_bytes = pio.to_image(fig, format="png", scale=1.5)
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.3),
                                     Inches(12.3), Inches(5.5))
            return True
        except Exception as e:
            logger.warning("Chart embed failed: %s", e)
            return False

    # ── ML Summary Slide ──────────────────────────────────────
    def add_ml_summary(self):
        ml = next((a for a in self.analyses if a.get("type") in ("ml", "classification", "regression", "clustering")), None)
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_INDIGO)

        _add_text(slide, "Machine Learning Summary", Inches(0.5), Inches(0.3),
                  Inches(10), Inches(0.6), font_size=26, bold=True, color=TEXT_WHITE)
        _add_divider(slide, Inches(1.0))

        content = self.exec_summary.get("model_performance", "")
        if not content and ml:
            results = ml.get("results", {})
            lines = [f"• {k.replace('_',' ').title()}: {v}"
                     for k, v in list(results.items())[:10]
                     if isinstance(v, (str, int, float))]
            content = "\n".join(lines)
        if not content:
            content = "No ML analysis data available. Run ML analysis to generate model performance metrics."

        _add_text(slide, content, Inches(0.5), Inches(1.3),
                  Inches(12.3), Inches(5.5), font_size=13, color=TEXT_MUTED)

    # ── Forecast Slide ─────────────────────────────────────────
    def add_forecast_slide(self):
        forecast = next((a for a in self.analyses if a.get("type") == "forecast"), None)
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_CYAN)

        _add_text(slide, "Time Series Forecast", Inches(0.5), Inches(0.3),
                  Inches(10), Inches(0.6), font_size=26, bold=True, color=TEXT_WHITE)
        _add_divider(slide, Inches(1.0), color=ACCENT_CYAN)

        if forecast:
            chart_added = self._try_embed_chart(slide, forecast)
            if not chart_added:
                res = forecast.get("results", {})
                lines = [
                    f"• Model: {res.get('model_type','N/A')}",
                    f"• Target Column: {res.get('value_column','N/A')}",
                    f"• Periods Forecasted: {res.get('periods', 0)}",
                ]
                _add_text(slide, "\n".join(lines), Inches(0.5), Inches(1.3),
                          Inches(12.3), Inches(5), font_size=14, color=TEXT_MUTED)
            # Insight from executive summary
            fi = self.exec_summary.get("forecast_insights", "")
            if fi:
                _add_text(slide, fi[:300], Inches(0.5), Inches(5.8),
                          Inches(12.3), Inches(1.2), font_size=12, color=TEXT_ACCENT)
        else:
            _add_text(slide, "No forecast data available. Run time series forecasting to populate this slide.",
                      Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5),
                      font_size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ── Anomaly Slide ─────────────────────────────────────────
    def add_anomaly_slide(self):
        anomaly = next((a for a in self.analyses if a.get("type") == "anomaly"), None)
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, RGBColor(0xEF, 0x44, 0x44))

        _add_text(slide, "Anomaly Detection", Inches(0.5), Inches(0.3),
                  Inches(10), Inches(0.6), font_size=26, bold=True, color=TEXT_WHITE)
        _add_divider(slide, Inches(1.0), color=RGBColor(0xEF, 0x44, 0x44))

        if anomaly:
            res = anomaly.get("results", {})
            total = res.get("total_anomalies", 0)
            method = res.get("method", "N/A")
            lines = [f"• Method: {method}", f"• Total Anomalies Detected: {total}"]
            by_col = res.get("anomalies_by_column", {})
            for col, info in list(by_col.items())[:6]:
                if info.get("count", 0) > 0:
                    lines.append(f"  ↳ {col}: {info['count']} anomalies ({info.get('percentage','?')}%)")
            _add_text(slide, "\n".join(lines), Inches(0.5), Inches(1.3),
                      Inches(8), Inches(4), font_size=13, color=TEXT_MUTED)

            ai_obs = self.exec_summary.get("anomaly_insights", "")
            if ai_obs:
                _add_text(slide, "AI Observations:", Inches(0.5), Inches(5.3),
                          Inches(12.3), Inches(0.35), font_size=12, bold=True, color=ACCENT_INDIGO)
                _add_text(slide, ai_obs[:350], Inches(0.5), Inches(5.7),
                          Inches(12.3), Inches(1.4), font_size=12, color=TEXT_ACCENT)
        else:
            _add_text(slide, "No anomaly data available. Run anomaly detection to populate this slide.",
                      Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5),
                      font_size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ── Recommendations Slide ─────────────────────────────────
    def add_recommendations(self):
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_PURPLE)

        _add_text(slide, "Strategic Recommendations", Inches(0.5), Inches(0.3),
                  Inches(10), Inches(0.6), font_size=26, bold=True, color=TEXT_WHITE)
        _add_divider(slide, Inches(1.0), color=ACCENT_PURPLE)

        content = self.exec_summary.get("recommendations", self.exec_summary.get("conclusion", ""))
        if not content:
            content = "Run a full analysis suite to generate AI-powered strategic recommendations."

        _add_text(slide, content[:900], Inches(0.5), Inches(1.3),
                  Inches(12.3), Inches(5.5), font_size=13, color=TEXT_MUTED)

    # ── Closing Slide ──────────────────────────────────────────
    def add_closing(self):
        slide = self.prs.slides.add_slide(self.blank)
        _set_slide_bg(slide, BG_DARK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_INDIGO)

        _add_text(slide, "DeepInsights", Inches(2), Inches(2.5),
                  Inches(9.33), Inches(1.2), font_size=48, bold=True,
                  color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        _add_text(slide, "Enterprise AI Analytics Platform",
                  Inches(2), Inches(3.7), Inches(9.33), Inches(0.6),
                  font_size=18, color=ACCENT_INDIGO, align=PP_ALIGN.CENTER)
        _add_divider(slide, Inches(4.5))
        _add_text(slide, "Generated by DeepInsights AI  •  Confidential",
                  Inches(2), Inches(4.7), Inches(9.33), Inches(0.5),
                  font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ── Build & Export ─────────────────────────────────────────
    def build(self) -> bytes:
        """Build all slides and return PPTX bytes."""
        self.add_cover()
        self.add_executive_summary()
        self.add_dataset_overview()
        self.add_analysis_slides()
        self.add_ml_summary()
        self.add_forecast_slide()
        self.add_anomaly_slide()
        self.add_recommendations()
        self.add_closing()

        out = io.BytesIO()
        self.prs.save(out)
        logger.info("PPTX report built: %d slides", len(self.prs.slides))
        return out.getvalue()


def generate_pptx(report_data: dict) -> bytes:
    """Public entry point: generate and return PPTX bytes."""
    return PptxBuilder(report_data).build()
